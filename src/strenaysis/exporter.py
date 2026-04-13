from __future__ import annotations

from io import BytesIO
from typing import Any

from docx import Document
from docx.shared import Inches
from pptx import Presentation


DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def build_export_bundle(payload: dict[str, Any], export_format: str) -> tuple[bytes, str, str]:
    fmt = str(export_format or "").strip().lower()
    if fmt == "docx":
        data = _build_docx(payload)
        return data, "strenaysis-workflow.docx", DOCX_MIME
    if fmt == "pptx":
        data = _build_pptx(payload)
        return data, "strenaysis-workflow.pptx", PPTX_MIME
    raise ValueError("Unsupported export format.")


def _build_docx(payload: dict[str, Any]) -> bytes:
    document = Document()
    document.add_heading("Strenaysis Workflow Review", level=0)
    document.add_paragraph(payload.get("problem", ""))

    document.add_heading("Problem Review", level=1)
    _add_labeled_paragraph(document, "Main question", payload.get("problem", ""))
    _add_labeled_paragraph(document, "Detailed context", payload.get("problem_details", "No detailed bucket added yet."))
    _add_labeled_paragraph(document, "Problem type", payload.get("problem_type", "Not set"))
    _add_labeled_paragraph(document, "Assessment", payload.get("assessment_title", "No assessment added yet."))
    _add_labeled_paragraph(document, "Recap", payload.get("assessment_recap", "No recap added yet."))

    document.add_heading("Roadmap Review", level=1)
    for index, node in enumerate(payload.get("nodes", []), start=1):
        build = node.get("build", {}) if isinstance(node, dict) else {}
        title = str(node.get("title", f"Node {index}")) if isinstance(node, dict) else f"Node {index}"
        document.add_heading(f"{index}. {title}", level=2)
        _add_labeled_paragraph(document, "Node purpose", node.get("why", ""))
        _add_labeled_paragraph(document, "Breakdown", node.get("breakdown", ""))
        _add_labeled_paragraph(document, "Execution summary", build.get("execution_summary", ""))
        _add_labeled_paragraph(document, "Problem parse", build.get("extracted_context", ""))

        workstreams = build.get("workstreams", [])
        if workstreams:
            document.add_paragraph("Working structure", style="Intense Quote")
            for workstream in workstreams:
                name = str(workstream.get("name", "")).strip()
                purpose = str(workstream.get("purpose", "")).strip()
                document.add_paragraph(f"{name}: {purpose}", style="List Bullet")

        actions = build.get("execution_items", [])
        if actions:
            document.add_paragraph("Action items", style="Intense Quote")
            for action in actions:
                paragraph = document.add_paragraph(style="List Bullet")
                paragraph.add_run(str(action.get("action", "")).strip()).bold = True
                details = [
                    f"Owner: {action.get('owner', 'Not set')}",
                    f"Approver: {action.get('approval', 'Not set')}",
                    f"Source: {action.get('source', 'Not set')}",
                    f"Artifact: {action.get('artifact', 'Not set')}",
                ]
                document.add_paragraph(" | ".join(details))

        output_sections = build.get("output_sections", {}) if isinstance(build, dict) else {}
        document.add_paragraph("Deck background context", style="Intense Quote")
        _add_labeled_paragraph(document, "Focus", output_sections.get("focus", ""))
        _add_labeled_paragraph(document, "Work to complete", output_sections.get("work_to_complete", ""))
        _add_labeled_paragraph(document, "Owners and sources", output_sections.get("owners_and_sources", ""))
        _add_labeled_paragraph(document, "Risks and handoff", output_sections.get("risks_and_handoff", ""))

    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _build_pptx(payload: dict[str, Any]) -> bytes:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = "Strenaysis Workflow Review"
    title_slide.placeholders[1].text = str(payload.get("problem", ""))

    overview_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    overview_slide.shapes.title.text = "Problem Review"
    overview_tf = overview_slide.placeholders[1].text_frame
    overview_tf.clear()
    for line in [
        f"Main question: {payload.get('problem', '')}",
        f"Problem type: {payload.get('problem_type', 'Not set')}",
        f"Assessment: {payload.get('assessment_title', 'No assessment added yet.')}",
        f"Action coverage: {sum(len((node.get('build', {}) or {}).get('execution_items', [])) for node in payload.get('nodes', []))} items",
    ]:
        p = overview_tf.add_paragraph() if overview_tf.text else overview_tf.paragraphs[0]
        p.text = line
        p.level = 0

    for index, node in enumerate(payload.get("nodes", []), start=1):
        build = node.get("build", {}) if isinstance(node, dict) else {}
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = f"{index}. {node.get('title', f'Node {index}')}"
        tf = slide.placeholders[1].text_frame
        tf.clear()
        bullets = [
            f"Purpose: {node.get('why', '')}",
            f"Execution summary: {build.get('execution_summary', '')}",
            f"Focus: {(build.get('output_sections', {}) or {}).get('focus', '')}",
            f"Work to complete: {(build.get('output_sections', {}) or {}).get('work_to_complete', '')}",
        ]
        actions = (build.get("execution_items", []) or [])[:2]
        for action in actions:
            bullets.append(
                f"Action: {action.get('action', '')} | Owner: {action.get('owner', 'Not set')} | Approver: {action.get('approval', 'Not set')}"
            )
        for idx, bullet in enumerate(bullets):
            paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _add_labeled_paragraph(document: Document, label: str, value: str) -> None:
    paragraph = document.add_paragraph()
    run = paragraph.add_run(f"{label}: ")
    run.bold = True
    paragraph.add_run(str(value or ""))
